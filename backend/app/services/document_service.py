"""
Document Service - Advanced NLP-based text extraction and preprocessing
Supports: PDF, DOCX, PPTX, TXT with up to 100MB file size
"""

import os
import re
import uuid
import string
from typing import Tuple, List, Optional, Dict
from pathlib import Path
from PyPDF2 import PdfReader
from app.config import UPLOADS_DIR

# Optional imports for extended format support
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not available. DOCX support disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. PPTX support disabled.")

# NLP libraries
try:
    import nltk
    from nltk.corpus import stopwords
    from nltk.tokenize import sent_tokenize, word_tokenize
    NLTK_AVAILABLE = True
except ImportError:
    NLTK_AVAILABLE = False
    print("Warning: nltk not available. Basic text processing will be used.")

try:
    import spacy
    SPACY_AVAILABLE = True
except ImportError:
    SPACY_AVAILABLE = False
    print("Warning: spacy not available. Entity extraction disabled.")


class NLPTextPreprocessor:
    """Text preprocessing for TTS - only cleans unpronouncenable characters"""
    
    def __init__(self):
        self._nlp = None
        self._initialized = False
    
    def _lazy_init(self):
        """Lazy initialization of NLP resources"""
        if self._initialized:
            return
        
        if SPACY_AVAILABLE:
            try:
                self._nlp = spacy.load("en_core_web_sm")
            except OSError:
                print("spaCy model not found. Entity extraction disabled.")
                self._nlp = None
        
        self._initialized = True
    
    def preprocess_for_tts(self, text: str, remove_stopwords: bool = False, optimize_tts: bool = True) -> str:
        """
        Preprocess text for TTS - only removes unpronouncenable characters.
        Keeps all meaningful words including stop words for natural speech.
        """
        # Clean unpronouncenable characters only
        text = self._clean_unpronounceables(text)
        
        # Optionally expand abbreviations for better pronunciation
        if optimize_tts:
            text = self._expand_abbreviations(text)
        
        # Fix whitespace and basic formatting
        text = self._normalize_whitespace(text)
        
        return text
    
    def _clean_unpronounceables(self, text: str) -> str:
        """Remove only characters that cannot be pronounced by TTS"""
        # Remove control characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
        
        # Remove URLs (can't pronounce)
        text = re.sub(r'https?://\S+|www\.\S+', '', text)
        
        # Remove email addresses (sounds bad in TTS)
        text = re.sub(r'\S+@\S+\.\S+', '', text)
        
        # Remove page numbers like "Page 1 of 10"
        text = re.sub(r'\bPage\s+\d+\s*(of\s+\d+)?\b', '', text, flags=re.IGNORECASE)
        
        # Remove standalone numbers that are just page numbers (lines with only numbers)
        text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)
        
        # Remove special symbols that can't be spoken (keep basic punctuation)
        # Keep: letters, numbers, spaces, and basic punctuation . , ! ? ; : ' " - ( )
        text = re.sub(r'[^\w\s.,!?;:\'"()\-–—\n]', ' ', text)
        
        # Normalize different dash types to regular hyphen
        text = re.sub(r'[–—]', '-', text)
        
        # Normalize quotes
        text = text.replace('"', '"').replace('"', '"').replace(''', "'").replace(''', "'")
        
        # Remove excessive punctuation (e.g., !!! -> !)
        text = re.sub(r'([!?.]){2,}', r'\1', text)
        
        # Remove bullet point markers
        text = re.sub(r'^[\s]*[•●○■□▪▫►▻◆◇★☆→]\s*', '', text, flags=re.MULTILINE)
        
        return text
    
    def _normalize_whitespace(self, text: str) -> str:
        """Normalize whitespace while preserving paragraph structure"""
        # Replace multiple spaces with single space
        text = re.sub(r'[ \t]+', ' ', text)
        
        # Replace multiple newlines with double newline (paragraph break)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Clean up spaces around newlines
        text = re.sub(r' *\n *', '\n', text)
        
        return text.strip()
    
    def _expand_abbreviations(self, text: str) -> str:
        """Expand common abbreviations for better TTS pronunciation"""
        abbreviations = {
            r'\bDr\.': 'Doctor',
            r'\bMr\.': 'Mister',
            r'\bMrs\.': 'Misses',
            r'\bMs\.': 'Miss',
            r'\bProf\.': 'Professor',
            r'\bvs\.': 'versus',
            r'\betc\.': 'etcetera',
            r'\be\.g\.': 'for example',
            r'\bi\.e\.': 'that is',
            r'\bNo\.': 'number',
            r'\bSt\.': 'Street',
            r'\bAve\.': 'Avenue',
            r'\bBlvd\.': 'Boulevard',
        }
        for abbr, expansion in abbreviations.items():
            text = re.sub(abbr, expansion, text, flags=re.IGNORECASE)
        return text
    
    def extract_entities(self, text: str) -> Optional[Dict[str, List[str]]]:
        """Extract named entities from text using spaCy"""
        self._lazy_init()
        
        if not self._nlp:
            return None
        
        # Limit text length for performance
        max_length = 100000
        if len(text) > max_length:
            text = text[:max_length]
        
        doc = self._nlp(text)
        
        entities = {
            'persons': [],
            'organizations': [],
            'locations': [],
            'dates': [],
            'key_terms': []
        }
        
        for ent in doc.ents:
            if ent.label_ == 'PERSON':
                entities['persons'].append(ent.text)
            elif ent.label_ == 'ORG':
                entities['organizations'].append(ent.text)
            elif ent.label_ in ['GPE', 'LOC']:
                entities['locations'].append(ent.text)
            elif ent.label_ == 'DATE':
                entities['dates'].append(ent.text)
        
        # Extract key noun phrases
        for chunk in doc.noun_chunks:
            if len(chunk.text.split()) > 1:
                entities['key_terms'].append(chunk.text)
        
        # Deduplicate
        for key in entities:
            entities[key] = list(set(entities[key]))[:20]
        
        return entities
    
    def get_text_stats(self, text: str) -> Dict:
        """Get text statistics"""
        words = text.split()
        sentences = sent_tokenize(text) if NLTK_AVAILABLE else text.split('.')
        
        return {
            'character_count': len(text),
            'word_count': len(words),
            'sentence_count': len([s for s in sentences if s.strip()]),
            'avg_word_length': sum(len(w) for w in words) / max(len(words), 1),
            'avg_sentence_length': len(words) / max(len(sentences), 1)
        }


class DocumentService:
    """Document processing service with multi-format support and NLP preprocessing"""
    
    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
    
    SUPPORTED_FORMATS = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'ppt': 'application/vnd.ms-powerpoint',
        'txt': 'text/plain'
    }
    
    def __init__(self):
        self._ensure_uploads_dir()
        self.preprocessor = NLPTextPreprocessor()
    
    def _ensure_uploads_dir(self):
        """Ensure uploads directory exists"""
        os.makedirs(UPLOADS_DIR, exist_ok=True)
    
    def save_file(self, file_content: bytes, original_filename: str) -> str:
        """Save uploaded file with unique name"""
        ext = Path(original_filename).suffix.lower()
        unique_filename = f"{uuid.uuid4()}{ext}"
        file_path = os.path.join(UPLOADS_DIR, unique_filename)
        
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        return file_path
    
    def extract_text(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from document based on file type"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx(file_path)
        elif ext == '.txt':
            return self.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    
    def extract_text_from_pdf(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from PDF file"""
        try:
            reader = PdfReader(file_path)
            text_parts = []
            
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            
            text = '\n\n'.join(text_parts)
            metadata = {
                'pages_count': len(reader.pages),
                'file_type': 'pdf'
            }
            
            return text, metadata
        except Exception as e:
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")
    
    def extract_text_from_docx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            raise ValueError("DOCX support not available. Install python-docx.")
        
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            
            text = '\n\n'.join(text_parts)
            metadata = {
                'paragraphs_count': len(doc.paragraphs),
                'tables_count': len(doc.tables),
                'file_type': 'docx'
            }
            
            return text, metadata
        except Exception as e:
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")
    
    def extract_text_from_pptx(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            raise ValueError("PPTX support not available. Install python-pptx.")
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
            
            text = '\n\n'.join(text_parts)
            metadata = {
                'slides_count': len(prs.slides),
                'file_type': 'pptx'
            }
            
            return text, metadata
        except Exception as e:
            raise ValueError(f"Failed to extract text from PPTX: {str(e)}")
    
    def extract_text_from_txt(self, file_path: str) -> Tuple[str, Dict]:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            
            lines = text.split('\n')
            metadata = {
                'lines_count': len(lines),
                'file_type': 'txt'
            }
            
            return text, metadata
        except Exception as e:
            raise ValueError(f"Failed to read text file: {str(e)}")
    
    def process_document(
        self,
        file_content: bytes,
        original_filename: str,
        apply_nlp: bool = True,
        remove_stopwords: bool = False,
        optimize_for_tts: bool = True
    ) -> Dict:
        """Complete document processing pipeline"""
        # Validate file size
        if len(file_content) > self.MAX_FILE_SIZE:
            raise ValueError(f"File too large. Maximum size is {self.MAX_FILE_SIZE / 1024 / 1024}MB")
        
        # Save file
        file_path = self.save_file(file_content, original_filename)
        
        try:
            # Extract text
            raw_text, metadata = self.extract_text(file_path)
            
            if not raw_text.strip():
                raise ValueError("No text could be extracted from the document")
            
            # Get original stats
            original_stats = self.preprocessor.get_text_stats(raw_text)
            
            # Apply NLP preprocessing
            if apply_nlp:
                processed_text = self.preprocessor.preprocess_for_tts(
                    raw_text,
                    remove_stopwords=remove_stopwords,
                    optimize_tts=optimize_for_tts
                )
                entities = self.preprocessor.extract_entities(raw_text)
            else:
                processed_text = raw_text
                entities = None
            
            # Get processed stats
            processed_stats = self.preprocessor.get_text_stats(processed_text)
            
            # Calculate reduction
            reduction_ratio = round(
                (1 - len(processed_text) / max(len(raw_text), 1)) * 100, 1
            ) if apply_nlp else 0
            
            result = {
                'text': processed_text,
                'raw_text': raw_text,
                'file_path': file_path,
                'filename': original_filename,
                'file_type': metadata.get('file_type', 'unknown'),
                'original_character_count': original_stats['character_count'],
                'original_word_count': original_stats['word_count'],
                'processed_character_count': processed_stats['character_count'],
                'processed_word_count': processed_stats['word_count'],
                'sentence_count': processed_stats['sentence_count'],
                'reduction_ratio': reduction_ratio,
                'entities': entities,
                **metadata
            }
            
            return result
        
        finally:
            # Clean up file after processing
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except:
                pass
    
    @classmethod
    def get_supported_formats(cls) -> Dict:
        """Return supported file formats and their availability"""
        return {
            'pdf': True,
            'docx': DOCX_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'ppt': PPTX_AVAILABLE,
            'txt': True,
            'max_file_size_mb': cls.MAX_FILE_SIZE / 1024 / 1024
        }
